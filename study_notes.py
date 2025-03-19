import os
import tempfile
import streamlit as st
import fitz  # PyMuPDF for PDFs
from pptx import Presentation
from dotenv import load_dotenv
import together
from fpdf import FPDF

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics



# Load API key
load_dotenv()

# Retrieve API key
together_api_key = os.getenv("TOGETHER_API_KEY")
# Set API key for Together
os.environ["TOGETHER_API_KEY"] = together_api_key
# together.api_key = os.getenv("TOGETHER_API_KEY")  # This is now optional but can be kept

# Function to extract text from PDF
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "\n".join([page.get_text("text") for page in doc])
    print(text)

    return text

# Function to extract text from PPTX
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

# Function to chunk text
def chunk_text(text, max_length=500):
    return [text[i: i + max_length] for i in range(0, len(text), max_length)]
import together

def generate_notes(text):
    prompt = f"""
    You are an AI that creates well-structured and concise study notes from academic materials.  
    Summarize the text below into **clear, structured, and well-organized** study notes.  
    **Instructions:**  
    - Use **headings** for different sections.  
    - Provide **bullet points** for key concepts.  
    - Use **examples** when needed.  
    - Summarize concisely without losing important details.  

    **Text:**  
    {text}

    The study notes should include:
    - A **brief summary** of key ideas.
    - **Bullet points** for important concepts.
    - **Definitions** of technical terms.
    - **Examples** where relevant.
    - Logical **organization** for better readability.

    **Study Notes Output:**  
    """

    response = together.Complete.create(
        model="meta-llama/Llama-2-13b-chat-hf",
        prompt=prompt,
        max_tokens=1000,
        temperature=0.3,
    )

    # Ensure response is valid
    if response and "output" in response:
        return response["output"]
    else:
        return "Error: Unable to generate summary."


    if response and hasattr(response, 'choices') and len(response.choices) > 0:
        return response.choices[0].text.strip()
    else:
        return "Error: Unable to generate summary."

# Function to save summary as PDF (with UTF-8 handling)
def save_summary_as_pdf(summary, output_path):
    c = canvas.Canvas(output_path, pagesize=letter)
    c.setFont("Helvetica", 12)  # Uses built-in font that supports UTF-8

    text_object = c.beginText(50, 750)
    text_object.setFont("Helvetica", 12)

    for line in summary.split("\n"):
        text_object.textLine(line)

    c.drawText(text_object)
    c.save()

# Streamlit UI
st.title("📚 AI-Powered Personalized Study Notes Generator")

uploaded_file = st.file_uploader("Upload a PDF or PPTX file", type=["pdf", "pptx"])

if uploaded_file:
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        temp_file.write(uploaded_file.read())
        temp_path = temp_file.name

    file_type = uploaded_file.name.split(".")[-1]

    if file_type == "pdf":
        extracted_text = extract_text_from_pdf(temp_path)
    elif file_type == "pptx":
        extracted_text = extract_text_from_pptx(temp_path)
    else:
        st.error("Unsupported file type.")
        extracted_text = None

    if extracted_text:
        if st.button("Generate Study Notes"):
            with st.spinner("Generating AI-powered notes..."):
                notes = generate_notes(extracted_text)

                output_pdf_path = "study_notes.pdf"
                save_summary_as_pdf(notes, output_pdf_path)

                st.success("Study notes generated successfully!")
                st.download_button(label="📥 Download Study Notes (PDF)", data=open(output_pdf_path, "rb"), file_name="study_notes.pdf", mime="application/pdf")
